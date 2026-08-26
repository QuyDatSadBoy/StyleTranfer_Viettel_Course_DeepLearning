"""
Sinh slide trình bày (.pptx) từ kết quả thực tế trong repo.

Chạy:  python slides/make_slides.py
Kết quả: slides/BaoCao_TangCuongDuLieuThoiTiet.pptx
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUT = ROOT / "slides" / "BaoCao_TangCuongDuLieuThoiTiet.pptx"

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(0x10, 0x1B, 0x33)
BLUE = RGBColor(0x2E, 0x6F, 0xD9)
TEAL = RGBColor(0x18, 0xA0, 0x99)
GREY = RGBColor(0x5A, 0x63, 0x75)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"


# --------------------------------------------------------------------------- #
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, spacing=1.15):
    """runs: list các (nội dung, size, bold, color) hoặc chuỗi."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    if isinstance(runs, str):
        runs = [(runs, 18, False, NAVY)]
    for i, item in enumerate(runs):
        content, size, bold, color = item if isinstance(item, tuple) else (item, 18, False, NAVY)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = content
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def header(slide, title: str, sub: str = "", n: int | None = None):
    rect(slide, 0, 0, W, Inches(1.05), NAVY)
    rect(slide, 0, Inches(1.05), W, Inches(0.055), TEAL)
    text(slide, Inches(0.6), Inches(0.16), Inches(11.4), Inches(0.9),
         [(title, 27, True, WHITE)] + ([(sub, 13, False, RGBColor(0xB6, 0xC2, 0xD9))] if sub else []))
    if n is not None:
        text(slide, Inches(12.2), Inches(6.95), Inches(0.9), Inches(0.4),
             [(str(n), 11, False, GREY)], align=PP_ALIGN.RIGHT)


def picture(slide, name: str, x, y, w=None, h=None):
    p = ASSETS / name
    if not p.exists():
        text(slide, x, y, Inches(6), Inches(0.5), [(f"(thiếu hình {name})", 12, False, GREY)])
        return None
    return slide.shapes.add_picture(str(p), x, y, width=w, height=h)


def fit_picture(slide, name: str, box_x, box_y, box_w, box_h):
    """Đặt ảnh vừa khít khung, giữ tỉ lệ, canh giữa."""
    from PIL import Image
    p = ASSETS / name
    if not p.exists():
        text(slide, box_x, box_y, Inches(6), Inches(0.5), [(f"(thiếu hình {name})", 12, False, GREY)])
        return
    iw, ih = Image.open(p).size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(p), Emu(int(box_x + (box_w - w) / 2)),
                             Emu(int(box_y + (box_h - h) / 2)), width=Emu(w))


def bullets(slide, x, y, w, items, size=17, gap=0.62):
    for i, it in enumerate(items):
        bold = it.startswith("**")
        content = it.replace("**", "")
        yy = Emu(int(y + Inches(gap) * i))
        rect(slide, x, Emu(int(yy + Inches(0.11))), Inches(0.11), Inches(0.11), TEAL)
        text(slide, Emu(int(x + Inches(0.3))), yy, w, Inches(0.6),
             [(content, size, bold, NAVY)])


def card(slide, x, y, w, h, title, body, accent=BLUE):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, Inches(0.07), h, accent)
    text(slide, Emu(int(x + Inches(0.28))), Emu(int(y + Inches(0.18))),
         Emu(int(w - Inches(0.5))), h,
         [(title, 16, True, NAVY), (body, 13, False, GREY)])


def table(slide, x, y, w, rows: list[list[str]], col_w: list[float],
          head_color=NAVY, size=13, row_h=0.42):
    """Bảng vẽ tay bằng shape để kiểm soát hoàn toàn kiểu dáng."""
    total = sum(col_w)
    for i, row in enumerate(rows):
        yy = Emu(int(y + Inches(row_h) * i))
        bg = head_color if i == 0 else (LIGHT if i % 2 else WHITE)
        rect(slide, x, yy, w, Inches(row_h), bg)
        cx = x
        for j, cell in enumerate(row):
            cw = Emu(int(w * col_w[j] / total))
            col = WHITE if i == 0 else NAVY
            text(slide, Emu(int(cx + Inches(0.12))), Emu(int(yy + Inches(0.07))),
                 cw, Inches(row_h), [(cell, size, i == 0, col)])
            cx = Emu(int(cx + cw))


# --------------------------------------------------------------------------- #
def load_results() -> tuple[dict, dict]:
    m = ROOT / "outputs" / "eval" / "metrics.json"
    d = ROOT / "experiments" / "detector" / "results.json"
    return (json.load(open(m)) if m.exists() else {},
            json.load(open(d)) if d.exists() else {})


def build() -> None:
    metrics, det = load_results()
    prs = new_deck()
    n = 0

    # ---------- 1. Bìa ---------- #
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, Inches(4.35), W, Inches(0.06), TEAL)
    text(s, Inches(0.9), Inches(2.05), Inches(11.7), Inches(2.2), [
        ("TĂNG CƯỜNG DỮ LIỆU ẢNH GIAO THÔNG", 33, True, WHITE),
        ("BẰNG CHUYỂN PHONG CÁCH THỜI TIẾT", 33, True, TEAL),
    ])
    text(s, Inches(0.9), Inches(4.62), Inches(11.7), Inches(1.0), [
        ("Đầu vào: 1 ảnh giao thông trời quang  +  1 ảnh tham chiếu thời tiết", 16, False, RGBColor(0xC8, 0xD2, 0xE4)),
        ("Đầu ra: ảnh giao thông dưới thời tiết xấu — giữ nguyên nhãn bounding box", 16, False, RGBColor(0xC8, 0xD2, 0xE4)),
    ])
    rect(s, Inches(0.9), Inches(5.95), Inches(0.06), Inches(1.0), TEAL)
    text(s, Inches(1.15), Inches(5.92), Inches(11.4), Inches(1.1), [
        ("NHÓM 6  ·  Bài tập cuối khoá — Deep Learning Ứng dụng", 15, True, TEAL),
        ("Hoàng Minh Đức  ·  Trần Quý Đạt  ·  Hoàng Trung Kiên  ·  Vũ Thuỳ Linh",
         14, False, RGBColor(0xC8, 0xD2, 0xE4)),
    ])

    # ---------- 2. Bài toán ---------- #
    n += 1
    s = blank(prs)
    header(s, "Bài toán & động lực", "Vì sao cần tăng cường dữ liệu thời tiết?", n)
    bullets(s, Inches(0.7), Inches(1.55), Inches(6.1), [
        "**Detector suy giảm mạnh khi thời tiết xấu.",
        "Dữ liệu huấn luyện lệch nặng về ảnh trời quang.",
        "BDD100K: ảnh mưa/tuyết ban ngày chỉ ~8%.",
        "**Thu thập ảnh thời tiết xấu rất tốn kém.",
        "Phải chờ đúng thời tiết + gán nhãn lại từ đầu.",
        "**Giải pháp: sinh ảnh từ dữ liệu đã có nhãn.",
        "Chi phí gán nhãn cho dữ liệu mới bằng 0.",
    ], size=16, gap=0.62)
    card(s, Inches(7.1), Inches(1.6), Inches(5.6), Inches(1.5), "ĐẦU VÀO",
         "1 ảnh giao thông điều kiện bình thường\n+ 1 ảnh tham chiếu thời tiết (mưa/tuyết/sương mù/mù khô/bão cát)", BLUE)
    card(s, Inches(7.1), Inches(3.3), Inches(5.6), Inches(1.5), "ĐẦU RA",
         "Chính ảnh giao thông đó dưới điều kiện thời tiết mong muốn,\nnhãn bounding box giữ nguyên 100%", TEAL)
    card(s, Inches(7.1), Inches(5.0), Inches(5.6), Inches(1.6), "RÀNG BUỘC",
         "· Mã nguồn mở  · Huấn luyện được trên 1 GPU\n· Suy luận thời gian thực  · Dễ giải thích", RGBColor(0xE4, 0x8A, 0x2E))

    # ---------- 3. Dữ liệu ---------- #
    n += 1
    s = blank(prs)
    header(s, "Dữ liệu", "Hai bộ dữ liệu công khai, tải tự động, không cần đăng ký", n)
    fit_picture(s, "fig_dataset.png", Inches(0.6), Inches(1.3), Inches(12.1), Inches(3.0))
    text(s, Inches(0.75), Inches(4.35), Inches(11.8), Inches(0.4),
         [("Biểu đồ = kho 10.000 ảnh có sẵn để chọn.   Bảng = số ảnh đã LỌC RA và thực sự dùng "
           "(chỉ lấy ban ngày, giới hạn số lượng cho nhẹ).", 13, False, GREY)])
    table(s, Inches(0.75), Inches(4.82), Inches(11.8), [
        ["Đã chọn ra từ kho", "Vai trò", "Số lượng", "Giấy phép"],
        ["BDD100K — trời quang, ban ngày", "Ảnh nội dung (đầu vào) + nhãn bounding box", "2.400 ảnh", "BSD-3"],
        ["BDD100K — mưa/tuyết, ban ngày", "300 ảnh tham chiếu  +  500 ảnh tập test THẬT", "800 ảnh", "BSD-3"],
        ["DAWN (Mendeley)", "Ảnh tham chiếu: sương mù / mù khô / mưa / tuyết / bão cát", "1.027 ảnh", "Chỉ nghiên cứu"],
    ], [3.0, 5.0, 1.7, 1.9], size=12)

    # ---------- 3b. Dữ liệu — ảnh mẫu ---------- #
    n += 1
    s = blank(prs)
    header(s, "Dữ liệu — ảnh mẫu & định dạng nhãn",
           "Hai loại ảnh đầu vào, và nhãn đi kèm mỗi ảnh nội dung", n)
    fit_picture(s, "fig_data_samples.jpg", Inches(0.5), Inches(1.28), Inches(12.3), Inches(3.5))

    card(s, Inches(0.7), Inches(4.95), Inches(5.9), Inches(1.95),
         "MỖI ẢNH NỘI DUNG CÓ 3 LOẠI NHÃN",
         "· weather   — clear / overcast / rainy / snowy / foggy …\n"
         "· timeofday — daytime / night / dawn-dusk\n"
         "· bounding box — 10 lớp (car, pedestrian, traffic sign …)\n"
         "   62.468 box trên 3.200 ảnh, trung bình ~20 box/ảnh", BLUE)

    rect(s, Inches(7.0), Inches(4.95), Inches(5.6), Inches(1.95), NAVY)
    text(s, Inches(7.25), Inches(5.1), Inches(5.2), Inches(1.8), [
        ("ĐỊNH DẠNG NHÃN (YOLO)", 15, True, TEAL),
        ("b1c66a42-6f7d68ca.txt", 13, False, RGBColor(0x9F, 0xB0, 0xC8)),
        ("2  0.482  0.531  0.061  0.088     → car", 13, False, WHITE),
        ("9  0.781  0.391  0.031  0.062     → traffic sign", 13, False, WHITE),
        ("lớp   x_tâm   y_tâm   rộng   cao   (đã chuẩn hoá 0–1)", 12, False, RGBColor(0x9F, 0xB0, 0xC8)),
    ], spacing=1.1)

    # ---------- 4. Tổng quan phương pháp ---------- #
    n += 1
    s = blank(prs)
    header(s, "Phương pháp — 3 khối nối tiếp", "Mỗi khối bù đúng một điểm yếu của khối trước", n)
    xs = [Inches(0.6), Inches(4.85), Inches(9.1)]
    titles = ["① AdaIN + Decoder", "② Guided Filter", "③ Phủ hạt (vật lý)"]
    bodies = [
        "AdaIN nhuộm đặc trưng VGG của ảnh gốc bằng mean/std của ảnh thời tiết.\n\n"
        "Decoder dựng đặc trưng đã nhuộm trở lại thành ảnh RGB — phần DUY NHẤT được huấn luyện.\n\n"
        "→ Đổi tông màu, độ sáng, độ mù.",
        "Dùng ảnh gốc làm ảnh dẫn hướng, lọc tuyến tính cục bộ q = a·I + b.\n\n"
        "Bán kính lớn (32) → a, b mượt → đầu ra = ẢNH GỐC nhân trường tương phản + lệch màu.\n\n"
        "→ Giữ 100% chi tiết ảnh chụp.",
        "Sinh vệt mưa (nhiễu + motion blur) và bông tuyết (3 lớp độ sâu), phủ theo chế độ screen.\n\n→ Thêm chi tiết cục bộ mà AdaIN không tạo ra được.",
    ]
    colors = [BLUE, TEAL, RGBColor(0xE4, 0x8A, 0x2E)]
    for x, t, b, c in zip(xs, titles, bodies, colors):
        rect(s, x, Inches(1.45), Inches(3.65), Inches(3.6), LIGHT)
        rect(s, x, Inches(1.45), Inches(3.65), Inches(0.09), c)
        text(s, Emu(int(x + Inches(0.25))), Inches(1.68), Inches(3.2), Inches(3.4),
             [(t, 17, True, NAVY), (b, 12, False, GREY)])
    text(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.6), [
        ("Vì sao nhãn bounding box vẫn dùng lại được?", 19, True, NAVY),
        ("Cả 3 khối chỉ thay đổi GIÁ TRỊ MÀU tại từng pixel — không dịch chuyển, không co giãn, "
         "không xoay vật thể. Vị trí chiếc xe trong ảnh đầu ra trùng khít vị trí trong ảnh gốc, "
         "nên toàn bộ nhãn cũ được sao chép nguyên vẹn. Chi phí gán nhãn cho dữ liệu mới = 0.", 15, False, GREY),
    ])

    # ---------- 5. AdaIN + Decoder ---------- #
    n += 1
    s = blank(prs)
    header(s, "① Khối học sâu: AdaIN + Decoder",
           "AdaIN trộn thống kê (0 tham số) · Decoder dựng lại ảnh (phần DUY NHẤT được học)", n)

    rect(s, Inches(0.7), Inches(1.28), Inches(11.9), Inches(1.0), NAVY)
    text(s, Inches(0.7), Inches(1.44), Inches(11.9), Inches(0.8),
         [("AdaIN(c, s)  =  σ(s) · ( c − μ(c) ) / σ(c)  +  μ(s)", 26, True, TEAL)],
         align=PP_ALIGN.CENTER)

    # -- cột trái: AdaIN --
    rect(s, Inches(0.7), Inches(2.5), Inches(6.0), Inches(3.55), LIGHT)
    rect(s, Inches(0.7), Inches(2.5), Inches(6.0), Inches(0.08), BLUE)
    text(s, Inches(0.95), Inches(2.66), Inches(5.6), Inches(3.3), [
        ("AdaIN — 0 THAM SỐ HỌC", 16, True, NAVY),
        ("c = đặc trưng VGG ảnh giao thông", 12, False, GREY),
        ("s = đặc trưng VGG ảnh thời tiết", 12, False, GREY),
        ("Bước 1 — ( c − μ(c) ) / σ(c)", 13, True, BLUE),
        ("     xoá phong cách gốc (tông trời quang),", 12, False, GREY),
        ("     GIỮ NGUYÊN cấu trúc không gian", 12, False, GREY),
        ("Bước 2 — nhân σ(s), cộng μ(s)", 13, True, BLUE),
        ("     'nhuộm' bằng thống kê ảnh thời tiết", 12, False, GREY),
        ("Chỉ đổi thống kê THEO KÊNH ⇒ bản đồ đặc trưng", 12, True, NAVY),
        ("không bị dịch chuyển ⇒ giữ nguyên bố cục.", 12, True, NAVY),
    ], spacing=0.98)

    # -- cột phải: Decoder --
    rect(s, Inches(7.0), Inches(2.5), Inches(5.6), Inches(3.55), LIGHT)
    rect(s, Inches(7.0), Inches(2.5), Inches(5.6), Inches(0.08), RGBColor(0xE4, 0x8A, 0x2E))
    text(s, Inches(7.25), Inches(2.66), Inches(5.2), Inches(3.3), [
        ("DECODER — 3,51 M THAM SỐ, PHẦN DUY NHẤT ĐƯỢC HỌC", 14, True, NAVY),
        ("Vì sao cần? VGG chỉ NÉN ảnh thành đặc trưng, không", 12, False, GREY),
        ("có mạng nào giải nén ngược lại ⇒ phải tự huấn luyện.", 12, False, GREY),
        ("Kiến trúc — phản chiếu ngược VGG-19", 13, True, RGBColor(0xE4, 0x8A, 0x2E)),
        ("     512 kênh @ 32×32   (1/8 độ phân giải)", 12, False, GREY),
        ("        ↓  9 lớp Conv 3×3 + ReLU", 12, False, GREY),
        ("        ↓  3 lần Upsample ×2", 12, False, GREY),
        ("     3 kênh @ 256×256   (ảnh RGB)", 12, False, GREY),
        ("· ReflectionPad thay zero-pad → không có viền đen", 12, False, GREY),
        ("· Upsample nearest thay ConvTranspose → không bị", 12, False, GREY),
        ("   vệt caro (checkerboard artifact)", 12, False, GREY),
    ], spacing=0.96)

    text(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.9), [
        ("Hệ quả quan trọng: AdaIN không chứa tham số nào, còn Decoder chỉ học cách "
         "\"dựng ảnh từ đặc trưng\" — không học riêng loại thời tiết nào.", 14, True, NAVY),
        ("⇒ Thêm loại thời tiết mới chỉ cần đưa ảnh tham chiếu vào, KHÔNG phải huấn luyện lại.",
         14, False, GREY),
    ])

    # ---------- 6. Tiền xử lý & huấn luyện ---------- #
    n += 1
    s = blank(prs)
    header(s, "Tiền xử lý dữ liệu & quy trình huấn luyện",
           "Chỉ Decoder được học — VGG-19 đóng băng hoàn toàn", n)

    rect(s, Inches(0.7), Inches(1.35), Inches(6.0), Inches(3.75), LIGHT)
    rect(s, Inches(0.7), Inches(1.35), Inches(6.0), Inches(0.08), BLUE)
    text(s, Inches(0.95), Inches(1.5), Inches(5.6), Inches(3.5), [
        ("TIỀN XỬ LÝ", 17, True, NAVY),
        ("① Lọc dữ liệu — chỉ giữ ảnh BAN NGÀY", 14, True, BLUE),
        ("     clear / partly cloudy  →  ảnh nội dung", 12, False, GREY),
        ("     rainy / snowy  →  ảnh tham chiếu + tập test", 12, False, GREY),
        ("② Xử lý ảnh trước khi đưa vào mạng", 14, True, BLUE),
        ("     Ảnh gốc 1280×720", 12, False, GREY),
        ("        ↓  resize cạnh ngắn về 320 px", 12, False, GREY),
        ("        ↓  cắt ngẫu nhiên 256×256  (tăng đa dạng)", 12, False, GREY),
        ("        ↓  lật ngang + chuẩn hoá về [0, 1]", 12, False, GREY),
        ("     →  Tensor 8×3×256×256", 12, True, NAVY),
    ], spacing=0.98)

    rect(s, Inches(7.0), Inches(1.35), Inches(5.6), Inches(3.75), LIGHT)
    rect(s, Inches(7.0), Inches(1.35), Inches(5.6), Inches(0.08), TEAL)
    text(s, Inches(7.25), Inches(1.5), Inches(5.2), Inches(3.5), [
        ("MỘT BƯỚC HUẤN LUYỆN", 17, True, NAVY),
        ("1.  Lấy 8 ảnh trời quang + 8 ảnh thời tiết ngẫu nhiên", 12, False, GREY),
        ("2.  Đưa cả hai qua VGG-19  →  đặc trưng c và s", 12, False, GREY),
        ("3.  t = AdaIN(c, s)          (trộn thống kê)", 12, False, GREY),
        ("4.  out = Decoder(t)        (dựng lại thành ảnh)", 12, False, GREY),
        ("5.  Đưa out qua VGG lần nữa  →  tính 3 loss", 12, False, GREY),
        ("6.  Cập nhật Decoder bằng Adam", 12, False, GREY),
        ("Ghép cặp NGẪU NHIÊN mỗi bước ⇒ decoder học tái", 12, True, NAVY),
        ("tạo cho MỌI phong cách ⇒ thời tiết mới chỉ cần đưa", 12, True, NAVY),
        ("ảnh tham chiếu vào, KHÔNG phải huấn luyện lại.", 12, True, NAVY),
    ], spacing=0.98)

    table(s, Inches(0.7), Inches(5.32), Inches(11.9), [
        ["Thành phần", "Giá trị", "Thành phần", "Giá trị"],
        ["Encoder", "VGG-19 pretrain, đóng băng — 0 tham số học", "Batch / số bước", "8 / 12.000"],
        ["Decoder", "3,51 M tham số — DUY NHẤT được học", "Optimizer", "Adam, lr 1e-4, decay 5e-5"],
        ["Dữ liệu", "2.160 ảnh nội dung × 1.327 ảnh tham chiếu", "Thời gian", "37 phút (RTX 5060 Ti)"],
    ], [1.6, 4.6, 1.9, 3.0], size=12, row_h=0.42)

    # ---------- 6b. Hàm mất mát ---------- #
    n += 1
    s = blank(prs)
    header(s, "Hàm mất mát", "Ba thành phần, mỗi thành phần giữ một thứ khác nhau", n)
    rect(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.95), NAVY)
    text(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.7),
         [("L  =  L_content  +  10 · L_style  +  1 · L_identity", 26, True, TEAL)],
         align=PP_ALIGN.CENTER)

    loss_cards = [
        ("① L_content", "MSE( VGG(ảnh_ra) , t )",
         "Ép ảnh ra có đặc trưng đúng bằng t = AdaIN(c,s). Decoder vẽ sai vị trí xe là loss tăng ngay.",
         "→ GIỮ BỐ CỤC CẢNH", BLUE),
        ("② L_style", "Σ MSE(μ, σ) qua 4 tầng VGG",
         "Chỉ so trung bình và độ lệch chuẩn, KHÔNG so từng điểm ảnh — vì phong cách nằm ở thống kê.",
         "→ GIỐNG TÔNG THỜI TIẾT", RGBColor(0xE4, 0x8A, 0x2E)),
        ("③ L_identity", "MSE( Decoder(VGG(c)) , c )",
         "Nếu style = content thì đầu ra PHẢI bằng đúng ảnh gốc. Điểm nhóm bổ sung so với bài báo gốc.",
         "→ GIỮ CẤU TRÚC, nhãn còn dùng được", TEAL),
    ]
    for x, (title, formula, why, tag, col) in zip(
            [Inches(0.7), Inches(4.85), Inches(9.0)], loss_cards):
        rect(s, x, Inches(2.5), Inches(3.6), Inches(2.6), LIGHT)
        rect(s, x, Inches(2.5), Inches(3.6), Inches(0.08), col)
        text(s, Emu(int(x + Inches(0.22))), Inches(2.68), Inches(3.2), Inches(2.4), [
            (title, 16, True, NAVY),
            (formula, 13, True, col),
            (why, 12, False, GREY),
            (tag, 12, True, NAVY),
        ], spacing=1.08)

    text(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.4),
         [("Đường cong huấn luyện thực tế — loss giảm mạnh 3.000 bước đầu rồi đi ngang; "
           "từ bước 10.000 chỉ còn giảm ~1% mỗi 1.000 bước ⇒ dừng ở 12.000 bước.", 13, False, GREY)])
    fit_picture(s, "fig_loss.png", Inches(0.7), Inches(5.72), Inches(11.9), Inches(1.5))

    # ---------- 7-10. Kết quả định tính ---------- #
    for title, sub, fig in (
        ("Kết quả — 5 loại thời tiết", "Cùng một ảnh gốc, đổi ảnh tham chiếu là đổi thời tiết", "fig_weather_types.jpg"),
        ("Kết quả — ví dụ đầu vào/đầu ra", "Trái: 2 đầu vào · Phải: ảnh sinh ra", "fig_examples.jpg"),
        ("Đóng góp của từng khối (ablation)", "Từ trái sang phải: thêm dần từng thành phần", "fig_ablation.jpg"),
    ):
        n += 1
        s = blank(prs)
        header(s, title, sub, n)
        fit_picture(s, fig, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.6))

    # ---------- 11. Nhãn được bảo toàn ---------- #
    n += 1
    s = blank(prs)
    header(s, "Bằng chứng: nhãn cũ vẫn khớp", "Cùng một bộ bounding box vẽ lên ảnh gốc và ảnh tăng cường", n)
    fit_picture(s, "fig_labels.jpg", Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.9))
    text(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.8),
         [("Không cần gán nhãn lại: mỗi ảnh sinh ra chỉ việc sao chép file nhãn của ảnh gốc.", 16, True, NAVY)])

    # ---------- 12. Định lượng ---------- #
    n += 1
    s = blank(prs)
    header(s, "Đánh giá định lượng", "Giữ nội dung (SSIM/EdgeRecall) và giống thật (FID)", n)
    r = metrics.get("results", {})
    fmt = lambda v, p=4: (f"{v:.{p}f}" if isinstance(v, (int, float)) and v == v and abs(v) != float("inf") else "—")  # noqa: E731
    rows = [["Phương pháp", "SSIM ↑", "PSNR ↑", "EdgeRec ↑", "FID ↓"]]
    if r:
        rows.append(["Không tăng cường (mốc domain gap)", "—", "—", "—", fmt(r.get("_no_augment", {}).get("fid"), 2)])
        for k, name in (("physics", "Baseline vật lý"), ("adain", "Chỉ AdaIN"),
                        ("ours", "AdaIN + GF + hạt (đề xuất)")):
            v = r.get(k, {})
            rows.append([name, fmt(v.get("ssim")), fmt(v.get("psnr"), 2),
                         fmt(v.get("edge_recall")), fmt(v.get("fid"), 2)])
    else:
        rows.append(["(chạy `python evaluate.py` để điền số liệu)", "", "", "", ""])
    table(s, Inches(0.7), Inches(1.35), Inches(11.9), rows, [4.2, 1.9, 1.9, 2.0, 1.9], row_h=0.46)
    table(s, Inches(0.7), Inches(4.0), Inches(11.9), [
        ["Chỉ số", "Cách tính", "Trả lời câu hỏi gì"],
        ["SSIM ↑", "So cấu trúc cục bộ giữa ảnh gốc và ảnh sinh ra", "Ảnh có giữ nội dung không?"],
        ["PSNR ↑", "10·log₁₀(255² / MSE) — sai lệch theo từng điểm ảnh", "Sai lệch nhiều hay ít?"],
        ["EdgeRecall ↑", "Tách biên Canny cả hai ảnh, đếm % biên GỐC còn tìm thấy",
         "Hình dáng vật thể còn nguyên? ⇒ nhãn còn dùng được?"],
        ["FID ↓", "Khoảng cách phân bố đặc trưng InceptionV3 với ảnh mưa/tuyết THẬT",
         "Ảnh sinh ra có giống ảnh chụp thật không?"],
    ], [1.7, 5.8, 4.4], size=12, row_h=0.44)
    text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.8), [
        ("Không phương pháp nào thắng tuyệt đối: bản đề xuất giữ nội dung tốt nhất (SSIM, PSNR), "
         "còn AdaIN thuần thắng FID vì lớp phủ hạt là dấu vết tổng hợp bị Inception phạt.", 13, False, GREY),
        ("⇒ Với ảnh dashcam nên giảm mật độ hạt — mưa thật qua kính chắn gió gần như không thấy vệt rời.",
         13, True, NAVY),
    ])

    # ---------- 13. Thí nghiệm YOLO ---------- #
    n += 1
    s = blank(prs)
    header(s, "Thí nghiệm quyết định: dữ liệu này có thực sự hữu ích?", "YOLOv8n — chỉ khác nhau ở tập huấn luyện", n)
    has_long_arm = "baseline_long" in det.get("models", {})
    line1 = ("A. Baseline: ảnh trời quang   ·   B. Baseline-long: cùng ảnh đó, gấp đôi epoch   ·   "
             "C. Augmented: ảnh trời quang + ảnh sinh ra") if has_long_arm else (
             "A. Baseline: chỉ ảnh trời quang   ·   C. Augmented: ảnh trời quang + ảnh do mô hình sinh ra")
    line2 = ("Nhánh B có ĐÚNG số bước cập nhật như C nhưng không có dữ liệu mới ⇒ so sánh công bằng là C − B. "
             if has_long_arm else "")
    text(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(1.3), [
        (line1, 15, True, NAVY),
        (line2 + "Cả hai đánh giá trên cùng tập test: ảnh mưa/tuyết THẬT của BDD100K, "
         "chưa từng dùng ở bất kỳ khâu nào khác.", 13, False, GREY),
    ])
    m = det.get("models", {})
    if m:
        has_long = "baseline_long" in m
        cols = ["baseline"] + (["baseline_long"] if has_long else []) + ["augmented"]
        vi_col = {"baseline": "A. Baseline", "baseline_long": "B. Baseline-long",
                  "augmented": "C. + Tăng cường"}
        rows = [["Tập kiểm tra"] + [vi_col[c] for c in cols] + ["C − " + ("B" if has_long else "A")]]
        for tag, vi in (("test_adverse", "Thời tiết xấu THẬT"), ("test_clear", "Trời quang (đối chứng)")):
            for k in ("mAP50", "mAP50-95"):
                ref = m["baseline_long"][tag][k] if has_long else m["baseline"][tag][k]
                d = m["augmented"][tag][k] - ref
                rows.append([f"{vi} · {k}"] + [f"{m[c][tag][k]:.4f}" for c in cols]
                            + [f"{d:+.4f} ({d/max(ref,1e-9)*100:+.1f}%)"])
        widths = [3.6] + [1.9] * len(cols) + [2.4]
        table(s, Inches(0.7), Inches(2.85), Inches(11.9), rows, widths, size=12, row_h=0.5)
    else:
        text(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(0.6),
             [("(chạy `python experiments/detector_experiment.py` để điền số liệu)", 15, False, GREY)])
    text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.9),
         [("Hàng 'trời quang' là đối chứng: xác nhận việc thêm dữ liệu tăng cường "
           "không làm mô hình kém đi trên điều kiện bình thường.", 14, False, GREY)])

    # ---------- 15. Kết luận & hướng phát triển ---------- #
    n += 1
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, W, Inches(1.0), RGBColor(0x0A, 0x12, 0x24))
    rect(s, 0, Inches(1.0), W, Inches(0.05), TEAL)
    text(s, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.7),
         [("KẾT LUẬN & HƯỚNG PHÁT TRIỂN", 30, True, WHITE)])

    rect(s, Inches(0.7), Inches(1.35), Inches(0.06), Inches(2.9), TEAL)
    text(s, Inches(0.95), Inches(1.3), Inches(11.6), Inches(3.0), [
        ("· Pipeline hoàn chỉnh sinh ảnh thời tiết xấu từ 2 ảnh đầu vào, đúng yêu cầu đề bài.", 16, False, RGBColor(0xD5, 0xDE, 0xEC)),
        ("· Mã nguồn mở; chỉ 3,51 M tham số được học, huấn luyện 37 phút trên 1 GPU tiêu dùng.", 16, False, RGBColor(0xD5, 0xDE, 0xEC)),
        ("· Nhãn bounding box tái sử dụng 100% — chi phí gán nhãn cho dữ liệu mới bằng 0.", 16, False, RGBColor(0xD5, 0xDE, 0xEC)),
        ("· Kiểm chứng bằng thí nghiệm thật: mAP50 trên ảnh thời tiết xấu thật tăng +6,2%.", 16, False, RGBColor(0xD5, 0xDE, 0xEC)),
        ("· Bàn giao mã nguồn chạy bằng 1 lệnh, web demo Gradio, báo cáo và slide.", 16, False, RGBColor(0xD5, 0xDE, 0xEC)),
    ], spacing=1.22)

    limits = ("· AdaIN chuyển thống kê TOÀN CỤC, chưa phân biệt vùng trời / mặt đường.\n"
              "· Không có bản đồ độ sâu thật; module vật lý giả thiết mặt đường phẳng.\n"
              "· Lớp phủ hạt làm giảm FID trên ảnh dashcam.\n"
              "· DAWN chỉ dùng cho nghiên cứu, không dùng thương mại.")
    future = ("· Thêm mặt nạ phân vùng trời / đường để đổi tông theo từng vùng.\n"
              "· Dùng mô hình ước lượng độ sâu đơn ảnh (Depth Anything).\n"
              "· Thay AdaIN bằng AdaAttN / WCT² để bám cấu trúc tốt hơn.\n"
              "· Thu thập ảnh tham chiếu thời tiết tại Việt Nam.")
    for x, title, body, col in ((Inches(0.7), "HẠN CHẾ", limits, RGBColor(0xE4, 0x5A, 0x5A)),
                                (Inches(6.95), "HƯỚNG PHÁT TRIỂN", future, TEAL)):
        rect(s, x, Inches(4.5), Inches(5.85), Inches(2.1), RGBColor(0x16, 0x22, 0x3C))
        rect(s, x, Inches(4.5), Inches(0.06), Inches(2.1), col)
        text(s, Emu(int(x + Inches(0.25))), Inches(4.65), Inches(5.4), Inches(1.95), [
            (title, 14, True, col),
            (body, 12, False, RGBColor(0xB6, 0xC2, 0xD9)),
        ], spacing=1.15)

    text(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5),
         [("Nhóm 6  —  Xin cảm ơn!", 20, True, TEAL)], align=PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"✓ Đã tạo slide: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} trang)")


if __name__ == "__main__":
    build()
